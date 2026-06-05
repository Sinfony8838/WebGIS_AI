"""Timeline generation service.

Parses lesson plan documents (.pptx, .pdf, .txt) and uses the LLM to
extract structured teaching stages into a timeline.
"""
from __future__ import annotations

import json
import re
from datetime import datetime, timezone
from typing import Any, Dict, List
from uuid import uuid4

from .minimax_client import LLMClient


TIMELINE_SYSTEM_PROMPT = """你是一个教学设计助手。用户会给你一份教案或教学设计文档的内容。
请从中提取教学流程，输出一个 JSON 数组，每个元素代表一个教学阶段。

每个元素的格式:
{
  "stage": "阶段名称(如: 导入、新授、练习、小结、作业布置)",
  "title": "具体标题",
  "description": "简要描述该阶段的教学活动",
  "durationMin": 预估时长(分钟, 整数)
}

要求:
1. 按教学时间顺序排列
2. 总时长应大致对应一节课(40-45分钟)
3. 如果文档中没有明确的时间分配，请根据内容合理估算
4. 至少包含3个阶段，最多不超过8个
5. 只输出 JSON 数组，不要输出其他内容
"""


class TimelineService:
    def __init__(self, llm_client: LLMClient):
        self.llm = llm_client

    def extract_text(self, filename: str, raw_bytes: bytes) -> str:
        lower = filename.lower()
        if lower.endswith(".pptx"):
            return self._extract_from_pptx(raw_bytes)
        if lower.endswith(".pdf"):
            return self._extract_from_pdf(raw_bytes)
        if lower.endswith(".txt") or lower.endswith(".md"):
            return raw_bytes.decode("utf-8", errors="replace")
        raise ValueError(f"不支持的文件格式: {filename}")

    def generate_timeline(
        self, filename: str, raw_bytes: bytes, project_id: str
    ) -> Dict[str, Any]:
        text = self.extract_text(filename, raw_bytes)
        if not text.strip():
            raise ValueError("文件内容为空")

        if len(text) > 8000:
            text = text[:8000]

        messages = [
            {"role": "system", "content": TIMELINE_SYSTEM_PROMPT},
            {"role": "user", "content": f"以下是教案内容:\n\n{text}"},
        ]

        raw_response = self.llm.chat_completion(messages, temperature=0.3)
        nodes_raw = self._parse_llm_json(raw_response)

        nodes: List[Dict[str, Any]] = []
        total_duration = 0
        for i, node in enumerate(nodes_raw):
            duration = int(node.get("durationMin", 5))
            total_duration += duration
            nodes.append({
                "id": str(uuid4()),
                "order": i,
                "stage": node.get("stage", f"阶段{i + 1}"),
                "title": node.get("title", ""),
                "description": node.get("description", ""),
                "durationMin": duration,
                "active": i == 0,
            })

        now = datetime.now(timezone.utc).isoformat()
        timeline = {
            "id": str(uuid4()),
            "project_id": project_id,
            "source_file_name": filename,
            "title": nodes_raw[0].get("title", "教学流程") if nodes_raw else "教学流程",
            "totalDurationMin": total_duration,
            "nodes": nodes,
            "created_at": now,
            "updated_at": now,
        }
        return {"status": "success", "timeline": timeline}

    # ------------------------------------------------------------------
    # Private helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_from_pptx(raw_bytes: bytes) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            raise ValueError("python-pptx is not installed")

        from io import BytesIO

        prs = Presentation(BytesIO(raw_bytes))
        texts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts)

    @staticmethod
    def _extract_from_pdf(raw_bytes: bytes) -> str:
        try:
            import fitz  # type: ignore[import-untyped]
        except ImportError:
            raise ValueError("PyMuPDF is not installed")

        doc = fitz.open(stream=raw_bytes, filetype="pdf")
        texts = [page.get_text() for page in doc]
        return "\n".join(texts)

    @staticmethod
    def _parse_llm_json(raw: str) -> List[Dict[str, Any]]:
        text = raw.strip()
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)
        parsed = json.loads(text)
        if isinstance(parsed, list):
            return parsed
        if isinstance(parsed, dict) and "nodes" in parsed:
            return parsed["nodes"]
        raise ValueError("LLM did not return a valid JSON array")
